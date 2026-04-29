import collections
import collections.abc
collections.Container = collections.abc.Container # Safe monkey-patch for older pptx versions on Python 3.10+

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Aesthetic settings
DARK_YELLOW = RGBColor(204, 153, 0)
BLACK = RGBColor(0, 0, 0)
FONT_NAME = 'Calibri'

def apply_text_style(run, font_size=18, bold=False, color=BLACK):
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

# Helpers for layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(title_slide_layout)
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]

title_1.text = "Software Maintenance Model Selection"
for p in title_1.text_frame.paragraphs:
    for run in p.runs:
         apply_text_style(run, font_size=40, bold=True, color=DARK_YELLOW)

subtitle_1.text = (
    "Project Graft (Browser Script Management Platform)\n"
    "Course: Software Engineering and Project Management - 21CSC303J\n\n"
    "Team Members:\n"
    "Product Owner: Vatsala Singh (RA2311003010993)\n"
    "Scrum Master: Sarthak Lal (RA2311003010991)\n"
    "Developer: Arian Bhattacharjee (RA2311003011001)\n"
    "End User: Priyanshu Bharadwaj (RA2311003010984)"
)
for p in subtitle_1.text_frame.paragraphs:
    for run in p.runs:
        apply_text_style(run, font_size=16)

# Create a generic function for content slides
def create_content_slide(title_text, content_text):
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title_text
    for p in title_shape.text_frame.paragraphs:
        for run in p.runs:
            apply_text_style(run, font_size=32, bold=True, color=DARK_YELLOW)
            
    tf = body_shape.text_frame
    tf.text = content_text
    
    # Format each run
    for p in tf.paragraphs:
        # Default spacing
        p.space_before = Pt(12)
        for run in p.runs:
            apply_text_style(run, font_size=20)
            
    return slide

# Slide 2
create_content_slide(
    "Suitable Maintenance Process Model",
    "• Chosen Model: Iterative Enhancement Model\n\n"
    "• Definition: This model assumes software is never truly 'finished.'.\n"
    "  It evolves through continuous, localized cycles of analysis, design,\n"
    "  coding, testing, and deployment.\n\n"
    "• This perfectly fits Project Graft's dynamic web ecosystem."
)

# Slide 3
create_content_slide(
    "Justification: Agile & Modern Architecture",
    "• Technical Stack: Graft uses a modern Next.js stack managed under an Agile framework.\n\n"
    "• Synergy: The Iterative Enhancement Model mirrors Agile Scrum sprints perfectly.\n\n"
    "• Application: Maintenance tasks, like resolving API recursion bugs or updating robust team features, are treated as standard iterative development tasks rather than separate, monolithic patches."
)

# Slide 4
create_content_slide(
    "Justification: Adapting to Target Web Environments",
    "• Core Purpose: Graft's primary function is to inject targeted client-side scripts into external, third-party websites.\n\n"
    "• Shifting Ecosystem: Target websites frequently update their DOM structures and layouts without warning.\n\n"
    "• Constant Adaptation: Graft must be maintained iteratively to continuously adapt and repair injection scripts as these external environments organically shift over time."
)

# Slide 5
create_content_slide(
    "Justification: Alignment with Deployment Pipeline",
    "• Natural Workflow: Graft natively features a 'Test -> Review -> Live' script deployment pipeline.\n\n"
    "• Model Alignment: This built-in pipeline matches the Iterative Enhancement Model’s core philosophy:\n"
    "   1. Analyze a localized defect\n"
    "   2. Code a solution in the dashboard Monaco editor\n"
    "   3. Isolate the test to the author's local Companion extension\n"
    "   4. Deploy the enhancement team-wide automatically."
)

# Slide 6
create_content_slide(
    "The Graft Iterative Maintenance Cycle",
    "• Step 1: Monitor Operations (Review Extension Delivery APIs)\n\n"
    "• Step 2: Analyze Bug Reports (Gather End User Feedback on injections)\n\n"
    "• Step 3: Iterative Enhancement (Fix/Enhance code via the Next.js Dashboard)\n\n"
    "• Step 4: Deploy (Automatic Sync to Team Companion Extensions)"
)

# Save the presentation
output_path = r"D:\Coding\graft\Graft_Maintenance_Model.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
